import os
import copy
import asyncio
import tempfile
import html
import uuid
import re
from aiogram import Bot, types
from aiogram.enums import ChatAction
from pptx import Presentation
from docxtpl import DocxTemplate

from config import PPT_TEMPLATE, DOCX_TEMPLATE
from database import get_test_paper
from utils import parse_raw_text, convert_to_pdf, format_docx_option


def safe_replace_text_in_paragraph(paragraph, old_key, new_val):
    """Template ke Exact Bold, Font, Color aur Size ko safe rakhte hue text replace karta hai"""
    if old_key not in paragraph.text:
        return

    for run in paragraph.runs:
        if old_key in run.text:
            run.text = run.text.replace(old_key, new_val)
            return

    full_text = paragraph.text.replace(old_key, new_val)
    if paragraph.runs:
        paragraph.runs[0].text = full_text
        for run in paragraph.runs[1:]:
            run.text = ""


async def generate_and_send(bot: Bot, chat_id: int, doc_id: str, gen_type: str):
    """MongoDB से डेटा निकालकर PPT/Test/Answer PDF जनरेट करके भेजता है"""
    msg = None
    output_file = None
    generated_pdf_path = None
    try:
        row = get_test_paper(doc_id)
        if not row:
            await bot.send_message(chat_id, "❌ <b>ID नहीं मिला!</b> कृपया सही ID दर्ज करें।")
            return

        topic = row["topic"]
        raw_text = row["raw_text"]
        parsed_qs = parse_raw_text(raw_text)
        
        try:
            debug_check_parsed_data(doc_id)
        except Exception as dbg_err:
            print(f"Debug function error: {dbg_err}")

        await bot.send_chat_action(chat_id=chat_id, action=ChatAction.TYPING)
        msg = await bot.send_message(chat_id, f"⏳ <b>{gen_type}</b> जनरेट हो रहा है, कृपया प्रतीक्षा करें...")

        temp_dir = tempfile.gettempdir()
        session_id = uuid.uuid4().hex[:6]
        
        loop = asyncio.get_running_loop()

        if gen_type == "PPT":
            if not os.path.exists(PPT_TEMPLATE):
                await msg.edit_text("❌ <b>Template Missing:</b> `template.pptx` नहीं मिला!")
                return
                
            output_file = os.path.join(temp_dir, f"temp_{doc_id}_{session_id}.pptx")
            prs = Presentation(PPT_TEMPLATE)
            
            base_slide = prs.slides[0]
            blank_layout = base_slide.slide_layout

            total_qs = len(parsed_qs)
            for _ in range(total_qs - 1):
                new_slide = prs.slides.add_slide(blank_layout)
                sp_tree = getattr(new_slide.shapes, '_spTree', getattr(new_slide.shapes, '_sptree', None))
                for shape in base_slide.shapes:
                    new_el = copy.deepcopy(shape.element)
                    if sp_tree is not None:
                        sp_tree.insert_element_before(new_el, 'p:extLst')

            for index, (slide, q) in enumerate(zip(prs.slides, parsed_qs), 1):
                cl_a = re.sub(r'✅|\*|\[ans\]', '', q['a'], flags=re.IGNORECASE).strip()
                cl_b = re.sub(r'✅|\*|\[ans\]', '', q['b'], flags=re.IGNORECASE).strip()
                cl_c = re.sub(r'✅|\*|\[ans\]', '', q['c'], flags=re.IGNORECASE).strip()
                cl_d = re.sub(r'✅|\*|\[ans\]', '', q['d'], flags=re.IGNORECASE).strip()
                
                replacements = {
                    '{{TOPIC}}': str(topic), 
                    '{{QUESTION}}': f"Q{index}. {q['text']}",
                    '{{OPTION_A}}': f"A) {cl_a}", 
                    '{{OPTION_B}}': f"B) {cl_b}",
                    '{{OPTION_C}}': f"C) {cl_c}", 
                    '{{OPTION_D}}': f"D) {cl_d}"
                }
                
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            for key, val in replacements.items():
                                safe_replace_text_in_paragraph(p, key, val)

            prs.save(output_file)
            generated_pdf_path = output_file.rsplit('.', 1)[0] + ".pdf"
            await loop.run_in_executor(None, convert_to_pdf, output_file, generated_pdf_path)

        else:
            if not os.path.exists(DOCX_TEMPLATE):
                await msg.edit_text("❌ <b>Template Missing:</b> `template.docx` नहीं मिला!")
                return
                
            output_file = os.path.join(temp_dir, f"temp_{doc_id}_{session_id}.docx")
            doc = DocxTemplate(DOCX_TEMPLATE)
            show_answers = (gen_type == "Answer Test PDF")
            
            formatted_qs = []
            for q in parsed_qs:
                formatted_qs.append({
                    'text': q['text'].strip(),
                    'opt_a': format_docx_option("(a)", q['a'], show_answers),
                    'opt_b': format_docx_option("(b)", q['b'], show_answers),
                    'opt_c': format_docx_option("(c)", q['c'], show_answers),
                    'opt_d': format_docx_option("(d)", q['d'], show_answers),
                })
            
            doc.render({'topic_name': topic, 'questions': formatted_qs})
            doc.save(output_file)
            generated_pdf_path = output_file.rsplit('.', 1)[0] + ".pdf"
            await loop.run_in_executor(None, convert_to_pdf, output_file, generated_pdf_path)

        if generated_pdf_path and os.path.exists(generated_pdf_path):
            await bot.send_chat_action(chat_id=chat_id, action=ChatAction.UPLOAD_DOCUMENT)
            await bot.send_document(
                chat_id, 
                types.FSInputFile(generated_pdf_path),
                caption=f"📄 आपका <b>{gen_type}</b> तैयार है!\n🆔 <b>ID:</b> <code>{doc_id}</code>"
            )
            if msg:
                await msg.delete()
        else:
            if msg:
                await msg.edit_text("❌ <b>PDF Generation Error:</b> PDF फाइल नहीं बन सकी।")

    except Exception as e:
        error_msg = f"❌ <b>Error:</b> {html.escape(str(e))}"
        if msg:
            await msg.edit_text(error_msg)
        else:
            await bot.send_message(chat_id, error_msg)

    finally:
        await asyncio.sleep(1)
        if output_file and os.path.exists(output_file): 
            try: os.remove(output_file)
            except: pass
        if generated_pdf_path and os.path.exists(generated_pdf_path): 
            try: os.remove(generated_pdf_path)
            except: pass


def debug_check_parsed_data(doc_id: str):
    """डेटाबेस और पार्सर का लाइव आउटपुट चेक करने के लिए डिबग फ़ंक्शन"""
    row = get_test_paper(doc_id)
    if not row:
        print(f"❌ ID {doc_id} डेटाबेस में नहीं मिला!")
        return
    
    raw_text = row.get("raw_text", "")
    parsed = parse_raw_text(raw_text)
    
    print("\n" + "="*50)
    print(f"🔍 DEBUG REPORT FOR DOC ID: {doc_id}")
    print("="*50)
    print("1. RAW TEXT FROM DATABASE:")
    print(raw_text)
    print("-" * 50)
    print("2. PARSED RESULT (dictionary):")
    for idx, q in enumerate(parsed, 1):
        print(f"  Q{idx}: {q.get('text')}")
        print(f"     A) {q.get('a')}")
        print(f"     B) {q.get('b')}")
        print(f"     C) {q.get('c')}")
        print(f"     D) {q.get('d')}")
    print("="*50 + "\n")
